#!/usr/bin/env python3
"""Generate a one-slide summary of the Alloy architecture experiments (mHC / KEEL).

    python scripts/make_experiments_slide.py   ->   docs/experiments.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
GREY = RGBColor(0x55, 0x55, 0x55)
LITE = RGBColor(0xEE, 0xF1, 0xF7)
RED = RGBColor(0xB0, 0x30, 0x30)

MHC = [
    ("dataset", "gdn3-334M-20B", "gdn3-mHC-20B"),
    ("mmlu", "0.2421", "0.2337"),
    ("arc-e", "0.5564", "0.5046"),
    ("arc-c", "0.2884", "0.2338"),
    ("hellaswag", "0.4277", "0.3279"),
    ("piqa", "0.6817", "0.6665"),
    ("avg", "0.4393", "0.3933"),
]
KEEL = [
    ("dataset", "qwen3-329M-20B", "qwen3-KEEL-20B"),
    ("mmlu", "0.2294", "0.2295"),
    ("arc-e", "0.3889", "0.2609"),
    ("arc-c", "0.2389", "0.2116"),
    ("hellaswag", "0.2877", "0.2587"),
    ("piqa", "0.556", "0.5223"),
    ("avg", "0.3402", "0.2966"),
]


def _txt(slide, left, top, width, height, text, size, *, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def _table(slide, data, left, top, width):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, left, top, width, Inches(2.9)).table
    gt.columns[0].width = Inches(1.55)
    gt.columns[1].width = Inches(1.95)
    gt.columns[2].width = Inches(1.95)
    for r, row in enumerate(data):
        header = r == 0
        avg = data[r][0] == "avg"
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = cell.margin_bottom = Pt(1)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(12 if header else 11)
            run.font.bold = header or avg
            # the "with-technique" column at the avg row is the key (lower) number -> red
            run.font.color.rgb = RED if (avg and c == 2) else (NAVY if header else GREY)
            if header:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif avg:
                cell.fill.solid(); cell.fill.fore_color.rgb = LITE
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _txt(s, Inches(0.5), Inches(0.28), Inches(12.4), Inches(0.6),
         "Efficient Architecture Experiments — mHC / KEEL", 26, bold=True)
    _txt(s, Inches(0.5), Inches(0.92), Inches(12.4), Inches(0.4),
         "Variants composed with Alloy, trained & evaluated on Ascend NPU.  "
         "Metric = acc (not acc_norm); avg = mean of mmlu / arc-e / arc-c / hellaswag / piqa "
         "(gsm8k ≈ 0 across runs, excluded).", 12, color=GREY)

    _txt(s, Inches(0.5), Inches(1.55), Inches(5.4), Inches(0.4),
         "mHC (hc_mult=2) on GDN3-334M · 20B tokens", 15, bold=True)
    _table(s, MHC, Inches(0.5), Inches(2.0), Inches(5.45))

    _txt(s, Inches(6.9), Inches(1.55), Inches(6.0), Inches(0.4),
         "KEEL (D=64, hq16, hkv2) on Qwen3-329M · 20B tokens", 15, bold=True)
    _table(s, KEEL, Inches(6.9), Inches(2.0), Inches(5.45))

    _txt(s, Inches(0.5), Inches(5.25), Inches(12.4), Inches(1.6),
         "Takeaways", 15, bold=True)
    body = s.shapes.add_textbox(Inches(0.5), Inches(5.65), Inches(12.4), Inches(1.5)).text_frame
    body.word_wrap = True
    for i, line in enumerate([
        "•  mHC: 0.393 vs 0.439 without mHC — no positive gain this round; investigating "
        "head config (D=64/hq16/hkv2 vs d=128/hq4/hkv4) and token budget.  (ongoing)",
        "•  KEEL: 0.297 vs 0.340 plain — no positive gain yet; deep-narrow (32L) and "
        "KEEL+deep-narrow variants still training.  (ongoing)",
        "•  Ref: arxiv.org/pdf/2601.19895",
    ]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(12); r.font.color.rgb = GREY

    out = Path(__file__).resolve().parent.parent / "docs" / "experiments.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
